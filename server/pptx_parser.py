#!/usr/bin/env python3
"""
Parse a Tony Bianco range review PPTX file.

Layout per slide:
- One heading text box (top-left): contains the LAST name + optional notes
  e.g. "MADDI – LINING – NO TOE PUFF / NO COUNTER – SIZE 11"
  The last name is the first word before any dash/em-dash.
  The heading is identified by cyan (#00FFFF) highlight OR being the topmost text box.

- Multiple style text boxes (one per style column):
  First paragraph = style name line (e.g. "MADDI $199.95" or "MOMA – UNLINED")
  Subsequent paragraphs = SKU lines (colour + leather)

Highlight colours (on <a:highlight> element in run XML):
  #FF00FF  magenta  → specked (sample here)
  #00FFFF  cyan     → specked_no_sample (specked but no sample yet)
  #FFFF00  yellow   → not_specked
  #FF0000  red      → cancelled
  #00FF00  green    → ignore (Anthony's notes)
  none             → carry_over

Output: JSON array to stdout, one object per slide:
  { last, styles: [ { style, skus: [ { colour, leather, status } ] } ] }
"""

import sys
import json
import re

def get_run_highlight(run):
    """Return the hex highlight colour of a run, or None."""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.find(f'{{{ns}}}rPr')
    if rPr is None:
        return None
    hl = rPr.find(f'{{{ns}}}highlight')
    if hl is None:
        return None
    srgb = hl.find(f'{{{ns}}}srgbClr')
    if srgb is None:
        return None
    val = srgb.get('val', '')
    return '#' + val.upper() if val else None

def get_para_highlight(para):
    """Return the highlight colour of the first highlighted run in a paragraph."""
    for run in para.runs:
        c = get_run_highlight(run)
        if c is not None:
            return c
    return None

def classify_highlight(hex_color):
    """Map hex highlight colour to SKU status."""
    if hex_color is None:
        return 'carry_over'
    h = hex_color.upper()
    if h == '#FF00FF':
        return 'specked'           # magenta = sample here
    if h == '#00FFFF':
        return 'specked_no_sample' # cyan = specked, no sample
    if h == '#FFFF00':
        return 'not_specked'       # yellow = not specked yet
    if h == '#FF0000':
        return 'cancelled'         # red = cancelled
    if h == '#00FF00':
        return 'ignore'            # green = Anthony's notes, ignore
    return 'carry_over'

def extract_last_name(heading_text):
    """Extract last name from heading like 'MADDI – LINING – NO TOE PUFF'."""
    # Split on dash, em-dash, en-dash, dollar
    parts = re.split(r'[\u2013\u2014\-\$]', heading_text)
    name = parts[0].strip()
    # Keep only letters and spaces
    name = re.sub(r'[^A-Za-z\s]', '', name).strip()
    return name.upper() if name else None

def extract_style_name(line):
    """
    Extract style name from a style header line like:
    'MADDI $199.95', 'MOMA – UNLINED', 'MAMZELLE –', 'MAXY $189.95'
    Returns the style name (first word, all caps) or None if not a style line.
    """
    line = line.strip()
    if not line:
        return None
    # Split on dash, em-dash, en-dash, dollar, space+digit
    core = re.split(r'[\u2013\u2014\-\$]|\s+\d', line)[0].strip()
    words = core.split()
    if not words:
        return None
    first = words[0]
    # Must be all uppercase letters, 2-15 chars
    if re.match(r'^[A-Z]{2,15}$', first):
        return first
    return None

def split_colour_leather(text):
    """
    Split 'BLACK NAPPA' → ('BLACK', 'NAPPA')
    Split 'BLACK NAPPA / BLACK PATENT T/C' → ('BLACK', 'NAPPA')
    For two-tone SKUs (e.g. VINTAGE / CHOC UNIT, PATENT / BLACK PATENT TC),
    we use ONLY the primary leather (before the /) so it matches the DB key.
    The secondary material (toe cap, unit, etc.) is a construction detail, not
    a separate leather in the dashboard.
    Returns (colour, leather) or (None, None) if unparseable.
    """
    # First clean up asterisks and trailing notes
    text = re.sub(r'\s*\*.*$', '', text).strip()

    # Split on slash for two-tone — use only the primary (first) part
    slash_parts = [p.strip() for p in text.split('/')]
    primary = slash_parts[0].strip()
    return _split_single(primary)

LEATHER_TYPES = [
    'NAPPA PATENT', 'PATENT TC', 'HI SHINE', 'NAPPA METALLIC',
    'NAPPA', 'SUEDE', 'PATENT', 'CROCO', 'VINTAGE', 'VENICE',
    'NUBUCK', 'KID', 'METALLIC', 'SPECKLE', 'CRINKLE', 'VELVET',
    'SATIN', 'GLITTER', 'MESH', 'FABRIC', 'LEATHER', 'CANVAS',
    'BROCADE', 'WOVEN', 'NYLON', 'SHINE', 'COMO', 'TC',
]

def _split_single(text):
    """Split a single-material string into (colour, leather)."""
    text = text.strip().upper()
    if not text:
        return None, None
    # Try two-word leather types first (longest match)
    words = text.split()
    for i in range(len(words) - 1, -1, -1):
        # Try two-word match
        if i > 0:
            two = words[i-1] + ' ' + words[i]
            if two in LEATHER_TYPES:
                colour = ' '.join(words[:i-1]).strip()
                return (colour if colour else None), two
        # Try one-word match
        if words[i] in LEATHER_TYPES:
            colour = ' '.join(words[:i]).strip()
            return (colour if colour else None), words[i]
    # No leather found
    return None, None

def is_note_line(text):
    """Return True if this line looks like a note/instruction, not a SKU."""
    text = text.strip()
    if not text:
        return True
    # Contains lowercase = likely a note
    if re.search(r'[a-z]', text):
        return True
    # Known note keywords
    note_keywords = ['HEEL', 'HEIGHT', 'LINING', 'TOE PUFF', 'COUNTER',
                     'SIZE', 'CM', 'MM', 'NOTE', 'NB:', 'SEE', 'REFER',
                     'CHECK', 'TBC', 'TBA', 'PENDING', 'CONFIRM',
                     'MICRO STRETCH', 'MICRO STRECH']
    for kw in note_keywords:
        if kw in text.upper():
            return True
    return False

# Words that look like colours but are actually notes/instructions
NOT_COLOURS = {
    'ADD', 'NO', 'TBC', 'TBA', 'SEE', 'NB', 'NOTE', 'CHECK',
    'PENDING', 'CONFIRM', 'REFER', 'ALSO', 'PLUS', 'NEW', 'OLD',
    'YES', 'OR', 'AND', 'THE', 'FOR', 'WITH', 'FROM',
}

def parse_sku_line(text, highlight):
    """Parse a SKU line and return {colour, leather, status} or None."""
    text = text.strip()
    if not text or is_note_line(text):
        return None

    status = classify_highlight(highlight)
    if status == 'ignore':
        return None

    colour, leather = split_colour_leather(text)
    if not colour or not leather:
        return None

    # Reject if colour is a known non-colour instruction word
    if colour.upper() in NOT_COLOURS:
        return None
    # Reject if colour contains question/instruction words (e.g. "SHOULD WE ADD BLACK")
    colour_words = colour.upper().split()
    question_words = {'SHOULD', 'COULD', 'WOULD', 'POSSIBLE', 'MAYBE', 'CONSIDER',
                      'WHAT', 'WHY', 'HOW', 'WHEN', 'WHERE', 'WHICH', 'IF', 'POSSIBLE'}
    if question_words & set(colour_words):
        return None
    # Reject if colour is more than 4 words (likely a note or instruction)
    if len(colour_words) > 4:
        return None

    return {'colour': colour, 'leather': leather, 'status': status}

def parse_slide(slide):
    """Parse a single slide and return {last, styles} or None."""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Get all text boxes with content, sorted top→bottom, left→right
    boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            boxes.append(shape)
    if not boxes:
        return None
    boxes.sort(key=lambda s: (s.top, s.left))

    # Find the heading box.
    # The heading is the topmost text box whose FIRST PARAGRAPH looks like a heading
    # (contains a dash/em-dash separator like "DAZIE – YSL HEEL") OR is the topmost box.
    # We deliberately do NOT use cyan highlight to find the heading, because style columns
    # can contain cyan-highlighted SKU lines (specked_no_sample) which would be misidentified
    # as the heading and then skipped during style parsing.
    last_name = None
    heading_box = None
    for box in boxes:
        first_para_text = ''
        for para in box.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                first_para_text = t
                break
        # A heading line contains a dash/em-dash separator (e.g. "DAZIE – YSL HEEL")
        # OR it is the topmost box (fallback)
        has_dash = bool(re.search(r'[\u2013\u2014\-]', first_para_text))
        if has_dash or heading_box is None:
            candidate_last = extract_last_name(first_para_text)
            if candidate_last:
                heading_box = box
                last_name = candidate_last
                if has_dash:
                    break  # found a proper heading with a dash separator
    if not last_name:
        return None
    # Clean last name: take only the first word
    last_name = last_name.split()[0] if last_name.split() else last_name

    # Parse each remaining box as a potential style column
    styles = []
    for box in boxes:
        if box is heading_box:
            continue
        paras = [(p.text.strip(), get_para_highlight(p))
                 for p in box.text_frame.paragraphs
                 if p.text.strip()]
        if not paras:
            continue

        # First paragraph should be the style name
        first_text, first_hl = paras[0]

        # Skip if it looks like a note (heel height etc.)
        if is_note_line(first_text):
            continue

        style_name = extract_style_name(first_text)
        if not style_name:
            continue

        # Skip if style name is a leather/colour word (false positive)
        false_positives = {'HEEL', 'SKIN', 'MILK', 'WEDGE', 'SOLE', 'UPPER',
                           'LINING', 'INSOLE', 'OUTSOLE', 'COUNTER', 'TOE'}
        if style_name in false_positives:
            continue

        skus = []
        for text, highlight in paras[1:]:
            sku = parse_sku_line(text, highlight)
            if sku:
                skus.append(sku)

        styles.append({'style': style_name, 'skus': skus})

    return {'last': last_name, 'styles': styles}

def parse_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    results = []
    for i, slide in enumerate(prs.slides):
        try:
            data = parse_slide(slide)
            if data and data['last']:
                results.append(data)
        except Exception as e:
            results.append({'error': str(e), 'slide': i + 1})
    return results

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({'error': 'Usage: pptx_parser.py <file.pptx>'}))
        sys.exit(1)
    filepath = sys.argv[1]
    try:
        data = parse_pptx(filepath)
        print(json.dumps(data, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({'error': str(e)}))
        sys.exit(1)
